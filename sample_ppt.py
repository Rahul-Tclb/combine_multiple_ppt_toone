import pptx
import os
from pptx import Presentation
from io import BytesIO


# Paths to your PowerPoint files



def combine_presentations(ppt_files, max_slides , output_file ):

    if os.path.exists(output_file):
        print(f"{output_file} exists and will be overwritten.")

    # Create a new empty presentation
    combined_ppt = Presentation()
    
    # Slide counter
    slide_count = 0
    
    # Loop through each PowerPoint file
    for ppt_file in ppt_files:
        # Open the existing presentation
        presentation = Presentation(ppt_file)
        
        # Copy slides until we reach the max limit
        for slide in presentation.slides:
            # Stop if the limit is reached
            if slide_count >= max_slides:
                break

            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            combined_ppt.slide_width = slide_width
            combined_ppt.slide_height = slide_height

            
            # Add slide layout to combined presentation
            slide_layout = combined_ppt.slide_layouts[0]
            new_slide = combined_ppt.slides.add_slide(slide_layout)
            
            
            # Copy contents of each slide shape (title, text, images, etc.)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Copy text content
                    textbox = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    textbox.text = shape.text

                elif shape.shape_type ==13:
                    image_stream = BytesIO(shape.image.blob)
                    new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
                
                elif shape.has_table:
                    table = shape.table
                    rows , cols = len(table.rows) , len(table.columns)
                    new_table = new_slide.shapes.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height).table
                    for row in range(rows):
                        for col in range(cols):
                            new_table.cell(row, col).text = table.cell(row, col).text
            
            slide_count += 1

    # Save the combined presentation
    combined_ppt.save("combined_presentation.pptx")
    print(f"Combined presentation created with {slide_count} slides in '{output_file}'.")

# List of PowerPoint files to combine
ppt_files = [
    r"C:\Users\pc\Rahul\Python\Adhock\file_1.pptx",
    r"C:\Users\pc\Rahul\Python\Adhock\file_2.pptx",
    r"C:\Users\pc\Rahul\Python\Adhock\file_3.pptx"
]

# Combine presentations with a limit of 15 slides
combine_presentations(ppt_files, max_slides = 5 , output_file="combined_presentation.pptx")


